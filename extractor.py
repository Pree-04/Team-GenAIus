import os
from docx import Document
import csv
from openpyxl import load_workbook
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
import json
from PyPDF2 import PdfReader
import cv2
import pytesseract
from pptx import Presentation

# Get the directory of the current script
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Specify the output file where all extracted data will be consolidated
output_file = os.path.join(BASE_DIR, "ExtractedRawData.txt")

# Function to append text to the output file
def append_to_output(data):
    with open(output_file, "a", encoding="utf-8") as f:
        f.write(data + "\n")

# DOCX extraction
def extract_docx(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".docx"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting text from DOCX file: {filename}\n")

            doc = Document(file_path)
            append_to_output("Paragraphs:")
            for para in doc.paragraphs:
                append_to_output(para.text)

            append_to_output("\nTables:")
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    append_to_output("\t".join(row_data))

            append_to_output("-" * 80)

# CSV extraction
def extract_csv(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".csv"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting data from CSV file: {filename}\n")

            with open(file_path, mode='r') as file:
                reader = csv.reader(file)
                for row in reader:
                    append_to_output(', '.join(row))
            append_to_output("-" * 80)

# Excel extraction
def extract_excel(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".xlsx"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting data from Excel file: {filename}\n")

            workbook = load_workbook(file_path, data_only=True)
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                append_to_output(f"Sheet: {sheet}\n")
                for row in worksheet.iter_rows(values_only=True):
                    append_to_output("\t".join([str(cell) if cell else '' for cell in row]))
            append_to_output("-" * 80)

# HTML text extraction
def extract_HTMLText(folder_path):
    # Check if the folder exists
    if not os.path.exists(folder_path):
        print(f"The folder {folder_path} does not exist.")
        return

    # Iterate through each file in the folder
    for filename in os.listdir(folder_path):
        if filename.endswith(".txt"):  # Process only TXT files
            file_path = os.path.join(folder_path, filename)
            print(f"Extracting text from: {filename}")

            # Open the TXT file and read its content
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                print(content)
            print("-" * 80)

# JSON extraction
def extract_json(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".json"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting data from JSON file: {filename}\n")

            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                append_to_output(json.dumps(data, indent=4))
            append_to_output("-" * 80)

# Markdown extraction
def extract_md(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".md"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting text from Markdown file: {filename}\n")

            with open(file_path, 'r', encoding='utf-8') as file:
                append_to_output(file.read())
            append_to_output("-" * 80)

# PDF extraction
def extract_pdf(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".pdf"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting text from PDF file: {filename}\n")

            with open(file_path, "rb") as file:
                reader = PdfReader(file)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text = page.extract_text()
                    append_to_output(f"Page {page_num + 1}:\n{text}")
            append_to_output("-" * 80)

# Specify the path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = "C:/Program Files/Tesseract-OCR/tesseract.exe"  # Update this if needed

def ocr_with_opencv(image_path):
    # Load the image using OpenCV
    image = cv2.imread(image_path)

    # Convert the image to grayscale for better OCR accuracy
    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

    # Use Tesseract to do OCR on the preprocessed image
    text = pytesseract.image_to_string(gray_image)
    return text

def extract_png(folder_path):
    # Iterate through each file in the folder
    for filename in os.listdir(folder_path):
        if filename.lower().endswith(".png"):  # Process only PNG files
            image_path = os.path.join(folder_path, filename)  # Create full path

            # Logging
            append_to_output(f"Extracting text from PNG file: {filename}\n")

            # Perform OCR using the updated logic
            extracted_text = ocr_with_opencv(image_path)

            # Append the extracted text to the output
            append_to_output(extracted_text)
            append_to_output("-" * 80)  # Separator for clarity

# PPTX extraction
def extract_pptx(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting text from PPTX file: {filename}\n")

            presentation = Presentation(file_path)
            for slide_number, slide in enumerate(presentation.slides):
                append_to_output(f"\nSlide {slide_number + 1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        append_to_output(f"Text: {shape.text}")
                    if shape.has_table:
                        append_to_output("Table found:")
                        table = shape.table
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            append_to_output("\t".join(row_data))
            append_to_output("-" * 80)

# TXT extraction
def extract_txt(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".txt"):
            file_path = os.path.join(folder_path, filename)
            append_to_output(f"Extracting text from TXT file: {filename}\n")

            with open(file_path, 'r', encoding='utf-8') as file:
                append_to_output(file.read())
            append_to_output("-" * 80)

# HTML URL extraction using Selenium
def extract_html_urls():
    options = Options()
    options.headless = True
    driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)

    urls = input("Enter the URLs separated by commas: ").split(',')

    for idx, url in enumerate(urls):
        url = url.strip()
        driver.get(url)
        html_content = driver.page_source
        append_to_output(f"HTML content extracted from URL {url}:\n{html_content}")
        append_to_output("-" * 80)

    driver.quit()
    append_to_output("HTML URL extraction complete.")

# Main function to call all extractors
def main():
    # Define folder paths relative to the base directory
    folder_paths = {
        "docx": os.path.join(BASE_DIR, "Downloads", "docx"),
        "csv": os.path.join(BASE_DIR, "Downloads", "csv"),
        "xlsx": os.path.join(BASE_DIR, "Downloads", "xlsx"),
        "json": os.path.join(BASE_DIR, "Downloads", "json"),
        "md": os.path.join(BASE_DIR, "Downloads", "md"),
        "pdf": os.path.join(BASE_DIR, "Downloads", "pdf"),
        "png": os.path.join(BASE_DIR, "Downloads", "png"),
        "pptx": os.path.join(BASE_DIR, "Downloads", "pptx"),
        "txt": os.path.join(BASE_DIR, "Downloads", "txt")
    }

    # Clear the output file before new extraction
    open(output_file, 'w').close()

    # Extract from each format
    extract_docx(folder_paths["docx"])
    extract_csv(folder_paths["csv"])
    extract_excel(folder_paths["xlsx"])
    extract_HTMLText(folder_paths["txt"])
    extract_json(folder_paths["json"])
    extract_md(folder_paths["md"])
    extract_pdf(folder_paths["pdf"])
    extract_png(folder_paths["png"])
    extract_pptx(folder_paths["pptx"])
    extract_txt(folder_paths["txt"])
    extract_html_urls()

if __name__ == "__main__":
    main()
